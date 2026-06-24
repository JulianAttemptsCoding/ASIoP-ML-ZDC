"""
make_presentation.py  —  ZDC Energy Reconstruction presentation
White background, matching color scheme of 20260501_ZDC_MC.pptx.pdf
Run: python scripts/make_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

BASE = r"C:\Users\Julia\OneDrive\Desktop\coding\ASIoP\Explore MC-sim data for ZDC"
OUT  = os.path.join(BASE, "plots", "ZDC_Energy_Reconstruction.pptx")
PLOT = {
    "slide5": os.path.join(BASE, "plots", "slide5_energy_dump.png"),
    "slide6": os.path.join(BASE, "plots", "slide6_erec_dist.png"),
    "slide7": os.path.join(BASE, "plots", "slide7_gamma_resolution.png"),
    "slide8": os.path.join(BASE, "plots", "slide8_neutron_resolution.png"),
}

# ── colours (matched to PDF) ──────────────────────────────────────────────
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK      = RGBColor(0x00, 0x00, 0x00)
C_DARK_BLUE  = RGBColor(0x1F, 0x38, 0x64)   # dark navy — divider line
C_MED_BLUE   = RGBColor(0x2E, 0x75, 0xB6)   # medium blue — accents / labels
C_LIGHT_BLUE = RGBColor(0xBD, 0xD7, 0xEE)   # pale blue — info boxes, squares
C_DECO_BLUE  = RGBColor(0x4E, 0x96, 0xC8)   # square decoration dark
C_GRAY       = RGBColor(0x59, 0x59, 0x59)   # body secondary text
C_DARK_GRAY  = RGBColor(0x26, 0x26, 0x26)   # body primary text

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

TOTAL = 9   # total slides for footer


# ── helpers ───────────────────────────────────────────────────────────────

def new_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_WHITE
    return s

def rect(s, x, y, w, h, fill, line=None, line_w=Pt(0)):
    from pptx.util import Inches
    shape = s.shapes.add_shape(1,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def tb(s, x, y, w, h, text, size, bold=False, color=C_BLACK,
       align=PP_ALIGN.LEFT, italic=False):
    tf = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf.text_frame.word_wrap = True
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tf

def mtb(s, x, y, w, h, lines, size=14, color=C_DARK_GRAY,
        align=PP_ALIGN.LEFT, space_before=2):
    """lines = list of (text, bold, color_or_None)"""
    tf = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf.text_frame.word_wrap = True
    for i, (text, bold, col) in enumerate(lines):
        p = tf.text_frame.paragraphs[0] if i == 0 else tf.text_frame.add_paragraph()
        p.alignment = align
        p.space_before = Pt(space_before)
        r = p.add_run()
        r.text = text
        r.font.size  = Pt(size)
        r.font.bold  = bold
        r.font.color.rgb = col if col else color
    return tf

def blue_box(s, x, y, w, h, lines, size=13):
    """Light-blue info box, black text, like PDF bullet boxes."""
    shape = s.shapes.add_shape(1,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_LIGHT_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    for i, (text, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(2)
        r = p.add_run()
        r.text = text
        r.font.size  = Pt(size)
        r.font.bold  = bold
        r.font.color.rgb = C_DARK_GRAY
    return shape

def deco_squares(s):
    """Two overlapping blue squares top-right, matching PDF."""
    rect(s, 12.50, 0.10, 0.45, 0.45, C_LIGHT_BLUE)
    rect(s, 12.75, 0.00, 0.45, 0.45, C_DECO_BLUE)

def divider(s, y=1.18):
    """Full-width dark-blue horizontal rule + small light-blue block, like PDF."""
    rect(s, 0.00, y,       13.33, 0.055, C_DARK_BLUE)
    rect(s, 0.30, y-0.005, 1.20,  0.065, C_LIGHT_BLUE)

def footer(s, slide_num, title_text="ZDC Energy Reconstruction"):
    """Footer: date left | title center | page right  (like PDF)."""
    rect(s, 0, 7.30, 13.33, 0.001, C_DARK_BLUE)
    tb(s, 0.2,  7.33, 2.5,  0.25, "2026/06/24",  size=10, color=C_GRAY)
    tb(s, 4.0,  7.33, 5.33, 0.25, title_text,    size=10, color=C_GRAY,
       align=PP_ALIGN.CENTER)
    tb(s, 11.8, 7.33, 1.3,  0.25,
       f"{slide_num}/{TOTAL}", size=10, color=C_GRAY, align=PP_ALIGN.RIGHT)

def content_header(s, title, slide_num):
    """Standard content slide header: large black title + divider."""
    deco_squares(s)
    tb(s, 0.30, 0.08, 12.5, 0.95, title, size=32, bold=False, color=C_BLACK)
    divider(s, y=1.05)
    footer(s, slide_num)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
deco_squares(s)

# Large bold title (centered-ish, upper half)
tb(s, 1.0, 1.6, 11.0, 1.8,
   "ZDC Energy Reconstruction",
   size=44, bold=True, color=C_BLACK, align=PP_ALIGN.CENTER)
tb(s, 1.0, 3.2, 11.0, 0.6,
   "Gamma (0.7 – 40 GeV)   |   Neutron (20 – 300 GeV)",
   size=20, bold=False, color=C_DARK_GRAY, align=PP_ALIGN.CENTER)

# Divider bar (same style as PDF title slide)
divider(s, y=4.05)

# Author / affiliation below divider
tb(s, 1.0, 4.25, 11.0, 0.5,
   "EIC  ·  Zero Degree Calorimeter  ·  Monte Carlo Study",
   size=18, bold=True, color=C_BLACK, align=PP_ALIGN.CENTER)
tb(s, 1.0, 4.85, 11.0, 0.45,
   "ASIoP  ·  Academia Sinica  ·  Taiwan Group",
   size=16, bold=False, color=C_DARK_GRAY, align=PP_ALIGN.CENTER)

footer(s, 1)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — BACKGROUND
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
content_header(s, "Background", 2)

mtb(s, 0.35, 1.25, 12.5, 0.5, [
    ("What is the EIC and ZDC?", True, C_MED_BLUE),
], size=18)

mtb(s, 0.35, 1.65, 12.5, 1.0, [
    ("The Electron-Ion Collider (EIC) collides electrons with protons and heavy nuclei to study "
     "the internal structure of matter. The Zero Degree Calorimeter (ZDC) sits at 0 degrees at "
     "the end of the beamline, catching neutral particles that fly straight through the bending "
     "magnets — primarily neutrons and photons (gammas).", False, None),
], size=14)

mtb(s, 0.35, 2.7, 12.5, 0.5, [
    ("Detector design", True, C_MED_BLUE),
], size=18)

blue_box(s, 0.35, 3.1, 5.9, 2.1, [
    ("ECAL  —  Electromagnetic Calorimeter", True),
    ("  Material:   LYSO crystal + SiPM", False),
    ("  Size:       20x20 cells, 3x3x7 cm/cell, 60x60 cm total", False),
    ("  Detects:    photons, electrons", False),
    ("  Radiation length:  7 cm ~ 6.5 X0 in Z", False),
], size=13)

blue_box(s, 6.55, 3.1, 6.4, 2.1, [
    ("HCAL  —  Hadronic Calorimeter (sampling)", True),
    ("  Material:   steel absorber + scintillator tile + SiPM", False),
    ("  Size:       64 layers, 65x60 cm face, 163 cm deep", False),
    ("  Detects:    neutrons, hadrons", False),
    ("  Layer:      1 steel + 1 scintillator per sampling unit", False),
], size=13)

mtb(s, 0.35, 5.35, 12.5, 0.5, [
    ("Goal of this study", True, C_MED_BLUE),
], size=18)

mtb(s, 0.35, 5.75, 12.5, 1.2, [
    ("Using Monte Carlo simulation with a 25-degree opening angle beam spread, determine how "
     "well the ZDC reconstructs the energy of incident gammas (0.7 – 40 GeV) and neutrons "
     "(20 – 300 GeV), and verify whether performance meets the EIC physics requirements.", False, None),
], size=14)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — RAW DATA & VARIABLE DEFINITIONS
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
content_header(s, "From Raw Data to Graph Axes", 3)

mtb(s, 0.35, 1.2, 5.9, 0.4, [("Step 1 — Simulation output (raw data)", True, C_MED_BLUE)], size=16)
blue_box(s, 0.35, 1.6, 5.9, 2.2, [
    ("E_beam   True beam energy (GeV)", True),
    ("         Set by simulation. What we want to reconstruct.", False),
    ("", False),
    ("E_ECAL   Energy deposited in ECAL layer (GeV)", True),
    ("         Measured detector signal.", False),
    ("", False),
    ("E_HCAL   Energy deposited in HCAL layer (GeV)", True),
    ("         Measured detector signal.", False),
], size=12)

mtb(s, 0.35, 3.95, 5.9, 0.4, [("Step 2 — Energy reconstruction", True, C_MED_BLUE)], size=16)
blue_box(s, 0.35, 4.35, 5.9, 1.85, [
    ("E_rec = p0 + p1*E_ECAL + p2*E_HCAL", True),
    ("", False),
    ("p0, p1, p2 fit by least-squares regression.", False),
    ("3 methods (m0 linear / m1 ratio / m2 quadratic).", False),
    ("2 weightings (w0 equal / w1 sqrt(E_beam)).", False),
], size=12)

mtb(s, 6.55, 1.2, 6.4, 0.4, [("Step 3 — Per-event ratio", True, C_MED_BLUE)], size=16)
blue_box(s, 6.55, 1.6, 6.4, 1.5, [
    ("r  =  E_rec / E_beam", True),
    ("", False),
    ("r = 1.0  ->  perfect reconstruction", False),
    ("r > 1.0  ->  overestimated energy", False),
    ("r < 1.0  ->  underestimated energy", False),
], size=12)

mtb(s, 6.55, 3.25, 6.4, 0.4, [("Step 4 — DCB fit on r distribution", True, C_MED_BLUE)], size=16)
blue_box(s, 6.55, 3.65, 6.4, 1.35, [
    ("mu    = mean  of fit  ->  BIAS value", True),
    ("sigma = width of fit  ->  RESOLUTION value", True),
    ("", False),
    ("Bell curve with power-law tails (DCB).", False),
], size=12)

mtb(s, 6.55, 5.15, 6.4, 0.4, [("Step 5 — What each graph axis is", True, C_MED_BLUE)], size=16)
blue_box(s, 6.55, 5.55, 6.4, 1.6, [
    ("Slide 5  x: E_beam     y: mean(E_sub/E_beam)", False),
    ("Slide 6  x: E_rec/E_beam    y: event count", False),
    ("Slide 7-8  x: E_beam", False),
    ("           y: sigma (resolution plot)", False),
    ("           y: mu    (bias plot)", False),
], size=12)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — METHOD
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
content_header(s, "Method", 4)

cols = [
    ("Why two detectors?",
     "Gammas deposit 20-60% of energy in ECAL, almost nothing in HCAL. "
     "Neutrons are the opposite — nearly all energy in HCAL. "
     "Using both signals together gives a complete picture.",
     0.35),
    ("Why regression?",
     "Raw signals are not equal to energy. There is leakage, dead material, "
     "and detector non-linearity. Regression learns the correction from simulation.",
     4.60),
    ("Why sqrt(E) weighting?",
     "Equal weighting lets high-energy events dominate the fit "
     "(bigger absolute residuals in GeV). "
     "sqrt(E_beam) weighting balances accuracy across the full energy range.",
     8.85),
]
for title, body, x in cols:
    mtb(s, x, 1.25, 4.1, 0.4, [(title, True, C_MED_BLUE)], size=16)
    mtb(s, x, 1.7,  4.1, 1.6, [(body,  False, None)], size=13)

mtb(s, 0.35, 3.4, 12.5, 0.4, [("Why DCB fit?", True, C_MED_BLUE)], size=16)
mtb(s, 0.35, 3.85, 12.5, 0.9, [
    ("Real shower distributions have heavier tails than a Gaussian — energy leaking out the sides "
     "creates a slow left tail; upward shower fluctuations create a slow right tail. "
     "The Double Crystal Ball uses a Gaussian core with power-law tails to match this.", False, None),
], size=13)

mtb(s, 0.35, 4.85, 12.5, 0.4, [("The requirement curve", True, C_MED_BLUE)], size=16)
blue_box(s, 0.35, 5.3, 5.9, 2.0, [
    ("sigma(E) = p0 / sqrt(E)  +  0.05", True),
    ("", False),
    ("p0/sqrt(E)  stochastic term: shower fluctuations scale as", False),
    ("            1/sqrt(E) from Poisson statistics.", False),
    ("0.05        constant floor: calibration, dead material,", False),
    ("            leakage — cannot improve with more energy.", False),
], size=12)

blue_box(s, 6.55, 5.3, 6.4, 2.0, [
    ("Gamma requirement:   p0 = 0.20 (new)   p0 = 0.35 (old)", False),
    ("Neutron requirement: p0 = 0.35 (new)   p0 = 0.50 (old)", False),
    ("", False),
    ("6 method/weighting combos tested:", False),
    ("  m0w0  m1w0  m2w0  (equal-weighted)", False),
    ("  m0w1  m1w1  m2w1  (sqrt(E)-weighted)", False),
], size=12)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ENERGY DUMP GRAPH
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
content_header(s, "Energy Dump", 5)

s.shapes.add_picture(PLOT["slide5"], Inches(0.2), Inches(1.2), width=Inches(8.9))

mtb(s, 9.3, 1.25, 3.8, 0.4, [("Axes", True, C_MED_BLUE)], size=15)
blue_box(s, 9.3, 1.65, 3.8, 1.55, [
    ("x-axis:  E_beam (GeV)", True),
    ("         True simulated beam energy.", False),
    ("", False),
    ("y-axis:  E_sub / E_beam", True),
    ("         Fraction of beam energy seen", False),
    ("         by one detector layer.", False),
], size=12)

mtb(s, 9.3, 3.35, 3.8, 0.4, [("Observations", True, C_MED_BLUE)], size=15)
blue_box(s, 9.3, 3.75, 3.8, 3.45, [
    ("Gamma (top row):", True),
    ("ECAL fraction 0.5 -> 0.15.", False),
    ("Falls with energy (saturation).", False),
    ("HCAL fraction ~10^-3 — tiny.", False),
    ("Gammas are ECAL particles.", False),
    ("", False),
    ("Neutron (bottom row):", True),
    ("ECAL fraction ~0 — near zero.", False),
    ("HCAL fraction 0.01 – 0.02.", False),
    ("Neutrons are HCAL particles.", False),
    ("", False),
    ("Neither detector alone is enough.", False),
], size=12)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — E_rec DISTRIBUTION
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
content_header(s, "E_rec / E_beam Distribution  —  1 GeV Gamma", 6)

s.shapes.add_picture(PLOT["slide6"], Inches(1.2), Inches(1.2), width=Inches(7.2))

mtb(s, 8.7, 1.25, 4.3, 0.4, [("Axes", True, C_MED_BLUE)], size=15)
blue_box(s, 8.7, 1.65, 4.3, 1.8, [
    ("x-axis:  r = E_rec / E_beam", True),
    ("         Reconstructed over true energy.", False),
    ("         r = 1.0 is perfect.", False),
    ("", False),
    ("y-axis:  Counts", True),
    ("         Events per r bin.", False),
    ("Red curve:  DCB fit", True),
], size=12)

mtb(s, 8.7, 3.6, 4.3, 0.4, [("Observations", True, C_MED_BLUE)], size=15)
blue_box(s, 8.7, 4.0, 4.3, 3.2, [
    ("Mean  = 1.01 -> bias < 1%", True),
    ("Near-perfect reconstruction.", False),
    ("", False),
    ("Sigma = 0.15 -> 15% resolution", True),
    ("Expected: 1/sqrt(1 GeV) term large.", False),
    ("", False),
    ("Chi2/NDF = 0.78 -> good fit", True),
    ("", False),
    ("Left tail:  energy leakage out", False),
    ("            detector sides/back.", False),
    ("Right tail: upward shower fluct.", False),
], size=12)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — GAMMA RESOLUTION & BIAS
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
content_header(s, "Gamma:  Resolution & Bias  (0.7 – 40 GeV)", 7)

s.shapes.add_picture(PLOT["slide7"], Inches(0.1), Inches(1.15), width=Inches(13.1))

blue_box(s, 0.20, 5.72, 6.3, 1.6, [
    ("Resolution (left panel):", True),
    ("y = sigma from DCB fit,  x = E_beam.", False),
    ("Dashed lines = requirement curves sigma = p0/sqrt(E) + 0.05.", False),
    ("All methods fall below the new requirement (p0 = 0.20).", False),
    ("Detector passes for gamma.", False),
], size=12)

blue_box(s, 6.75, 5.72, 6.3, 1.6, [
    ("Bias (centre panel):", True),
    ("y = mu from DCB fit,  ideal = 1.0.", False),
    ("All methods within 0.90 – 1.06 across full energy range.", False),
    ("Small dip at 0.7 GeV for quadratic equal-weighted (m2 w0).", False),
    ("Gamma bias is well-controlled for all methods.", False),
], size=12)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — NEUTRON RESOLUTION & BIAS
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
content_header(s, "Neutron:  Resolution & Bias  (20 – 300 GeV)", 8)

s.shapes.add_picture(PLOT["slide8"], Inches(0.1), Inches(1.15), width=Inches(13.1))

blue_box(s, 0.20, 5.72, 6.3, 1.6, [
    ("Resolution (left panel):", True),
    ("y = sigma from DCB fit,  x = E_beam.", False),
    ("sqrt(E)-weighted methods (open markers) meet the new requirement", False),
    ("(p0 = 0.35) at most energies.", False),
    ("Equal-weighted methods are worse, especially below 50 GeV.", False),
], size=12)

blue_box(s, 6.75, 5.72, 6.3, 1.6, [
    ("Bias (centre panel):", True),
    ("Equal-weighted: mu up to 2 – 3.5 at 20 GeV.", False),
    ("Regression dominated by high-E events, fails at low energy.", False),
    ("sqrt(E)-weighted: mu = 1.1 – 1.3 at 20 GeV. Far more controlled.", False),
    ("sqrt(E) weighting is the recommended method for neutrons.", False),
], size=12)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — OBSERVATIONS & OUTLOOK
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
content_header(s, "Observations & Outlook", 9)

mtb(s, 0.35, 1.2, 12.5, 0.4, [("Observations", True, C_MED_BLUE)], size=17)
blue_box(s, 0.35, 1.62, 5.9, 2.55, [
    ("Gamma", True),
    ("All 6 methods meet the new resolution requirement", False),
    ("(20%/sqrt(E) + 5%) across 0.7 – 40 GeV.", False),
    ("Bias controlled within +/- 6% for all methods.", False),
    ("Linear and quadratic perform similarly.", False),
    ("Ratio method slightly worse at low energy.", False),
], size=12)

blue_box(s, 6.55, 1.62, 6.4, 2.55, [
    ("Neutron", True),
    ("sqrt(E)-weighted methods are clearly superior.", False),
    ("Equal-weighted bias reaches 3.5x at 20 GeV — unreliable.", False),
    ("sqrt(E)-weighted bias stays below 1.35x across all energies.", False),
    ("Resolution close to new requirement (35%/sqrt(E) + 5%);", False),
    ("may need tuning at low energies.", False),
], size=12)

mtb(s, 0.35, 4.35, 12.5, 0.4, [("Outlook", True, C_MED_BLUE)], size=17)

outlook = [
    ("Optimise neutron at low energy",
     "20–50 GeV range shows largest bias and worst resolution. "
     "Consider energy-dependent regression or neural-network approach.",
     0.35),
    ("Validate with full GEANT4 simulation",
     "Current MC is idealized geometry. Full simulation with beam backgrounds, "
     "crossing angle, and realistic digitization will stress-test results.",
     4.60),
    ("Extend to ML regression",
     "Linear regression is a baseline. Networks using ECAL+HCAL hit maps "
     "(not just sums) could reduce sigma and bias further.",
     8.85),
]
for title, body, x in outlook:
    blue_box(s, x, 4.78, 4.1, 2.55, [
        (title, True),
        ("", False),
        (body, False),
    ], size=12)


# ── save ─────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"[OK] Saved: {OUT}")
